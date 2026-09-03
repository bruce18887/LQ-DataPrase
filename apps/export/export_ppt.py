"""PPTX builder for batch charts export.

Extracted from views.py ``_batch_charts_pptx``.  Uses matplotlib for
histogram generation and python-pptx for slide assembly.
"""

import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches

from apps.analysis.services.statistics import (
    compute_cpk, compute_range_statistics, get_1d_from, filter_finite,
)
from apps.analysis.services.statistics.helpers import normal_pdf_curve
from apps.analysis.services.statistics.kde import GaussianKDE
from apps.export.histogram_grid import build_histogram_grid, finite_or_none


def build_batch_charts_pptx(datafile, df, metadata, params,
                            show_limit=True, show_3sigma=False,
                            show_4sigma=False, show_6sigma=True,
                            show_normal=False, show_kde=False):
    """Build batch charts PPTX with histogram slides.

    For each parameter in *params* a matplotlib histogram is rendered and
    embedded into a blank PPTX slide.  The overlay switches mirror the xlsx
    branch (``build_batch_charts_xlsx_with_charts``) so the same chart config
    produces the same picture in both formats (缺陷 #6: pptx 曾一个开关都不收).

    Parameters
    ----------
    datafile : DataFile
        Used only for filename generation (by the caller).
    df : pd.DataFrame
    metadata : dict
    params : list of str
        Parameter (column) names to chart.
    show_limit, show_3sigma, show_4sigma, show_6sigma, show_normal, show_kde : bool
        与前端图表配置一一对应的叠加开关。

    Returns
    -------
    bytes
        Serialized PPTX content.
    """
    prs = Presentation()
    # blank layout
    blank_layout = prs.slide_layouts[6]

    # Chinese font support
    plt.rcParams['font.sans-serif'] = ['SimHei', 'Microsoft YaHei', 'DejaVu Sans']
    plt.rcParams['axes.unicode_minus'] = False

    for param in params:
        if param not in df.columns:
            continue
        data_series = filter_finite(get_1d_from(df, param))
        if len(data_series) == 0:
            continue

        stats = compute_range_statistics(data_series, metadata, param)
        # 限值可能为 None（parse_limit_string 新语义）或非有限：统一收敛后再用，
        # compute_cpk 对 None 返回 0.0，绘图侧则跳过标记线（不画幻影 LSL/USL）
        rdl_min = finite_or_none(stats['rdl'][0])
        rdl_max = finite_or_none(stats['rdl'][1])
        cpk_result = compute_cpk(stats['mean'], stats['std'], rdl_min, rdl_max)
        cpk_val = cpk_result['cpk']
        cpk_level = cpk_result.get('cpk_color', 'gray')

        # Create matplotlib chart
        fig, ax = plt.subplots(figsize=(8, 4.5))
        values = data_series.to_numpy(dtype=float)
        # 分箱与屏幕/xlsx 导出同一套网格（缺陷 #4/#5）：27 边界 = 25 内边界
        # + ±inf 兜底，超范围值计入首/尾 bin；限值退化时回退数据范围。
        edges, centers, gap = build_histogram_grid(rdl_min, rdl_max, values)
        counts, _ = np.histogram(values, bins=edges)
        ax.bar(centers, counts, width=gap * 0.9, color='#1E88E5',
               edgecolor='white', alpha=0.85, label='数据分布')
        ax.set_xlim(float(centers[0]), float(centers[-1]))

        if show_limit and rdl_min is not None:
            ax.axvline(rdl_min, color='#C62828', linestyle='--', linewidth=2, label='LSL')
        if show_limit and rdl_max is not None:
            ax.axvline(rdl_max, color='#C62828', linestyle='--', linewidth=2, label='USL')

        peak = int(counts.max()) if len(counts) else 0
        for sigma, flag in ((3, show_3sigma), (4, show_4sigma), (6, show_6sigma)):
            band = stats.get(f's{sigma}')
            if flag and band and stats['std'] > 0:
                ax.axvline(band[0], color='#E65100', linestyle=':', linewidth=1.5,
                           label=f'{sigma}σL')
                ax.axvline(band[1], color='#E65100', linestyle=':', linewidth=1.5,
                           label=f'{sigma}σU')

        if show_normal and stats['std'] > 0 and peak > 0:
            # 公式单一来源 normal_pdf_curve（与屏幕 / xlsx 导出同源），此处只做
            # max-normalize 到频数轴
            curve = normal_pdf_curve(stats['mean'], stats['std'],
                                     float(centers[0]), float(centers[-1]))
            if curve:
                x_pdf = np.array([x for x, _ in curve])
                y_pdf = np.array([y for _, y in curve])
                if np.max(y_pdf) > 0:
                    ax.plot(x_pdf, y_pdf / np.max(y_pdf) * peak,
                            color='#F57F17', linewidth=2, label='正态分布')

        if show_kde and peak > 0:
            # Best-effort：退化数据拟合失败就省略曲线，不让导出整体失败
            try:
                if len(values) >= 3 and np.ptp(values) > 0:
                    kde = GaussianKDE(values, bw_method='silverman')
                    x_kde = np.linspace(float(centers[0]), float(centers[-1]), 200)
                    y_kde = kde(x_kde)
                    if np.max(y_kde) > 0:
                        ax.plot(x_kde, y_kde / np.max(y_kde) * peak,
                                color='#7B1FA2', linewidth=2, label='KDE曲线')
            except Exception:  # noqa: BLE001
                pass

        ax.set_title(
            f'{param}  |  CPK={cpk_val:.4f} ({cpk_level})  |  N={len(data_series)}',
            fontsize=11,
        )
        ax.set_xlabel(stats.get('unit', ''))
        ax.set_ylabel('Frequency')
        handles, labels = ax.get_legend_handles_labels()
        if handles:
            ax.legend(handles, labels, loc='upper right', fontsize=8)
        fig.tight_layout()

        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=120)
        buf.seek(0)
        plt.close(fig)

        slide = prs.slides.add_slide(blank_layout)
        left = Inches(0.5)
        top = Inches(0.5)
        slide.shapes.add_picture(buf, left, top, width=Inches(9), height=Inches(5.5))

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
