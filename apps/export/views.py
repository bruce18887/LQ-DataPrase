import os, io, json
import pandas as pd
import numpy as np
from django.conf import settings
from django.shortcuts import get_object_or_404
from django.http import FileResponse
from rest_framework import viewsets, status
from rest_framework.decorators import action
from rest_framework.response import Response
from rest_framework.permissions import IsAuthenticated
from rest_framework.parsers import JSONParser
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from apps.datafiles.models import DataFile
from apps.datafiles.parsers import get_parser
from apps.analysis.services.statistics import (
    detect_fail_data, compute_cpk, compute_range_statistics,
    get_columns_with_limits, get_1d_from, get_site_column,
    get_bin_column_name, parse_limit_string, compute_site_stats,
)

class ExportViewSet(viewsets.GenericViewSet):
    permission_classes = [IsAuthenticated]

    @action(detail=False, methods=['post'])
    def to_excel(self, request):
        file_id = request.data.get('file_id')
        passfail = request.data.get('passfail', '全部')
        site_filter = request.data.get('site_filter', '全部')

        datafile = get_object_or_404(DataFile, pk=file_id, owner=request.user)
        parser = get_parser(datafile.format_type)
        df, metadata = parser.parse(datafile.file_path)
        if df is None:
            return Response({'error': 'parse_failed'}, status=400)

        export_df = df.copy()

        if site_filter != '全部':
            site_col = get_site_column(df)
            if site_col:
                export_df = export_df[export_df[site_col].astype(str) == str(site_filter)]

        if passfail != '全部':
            fail_indices, _, _ = detect_fail_data(export_df, metadata)
            fail_set = set(fail_indices)
            if passfail == 'Fail':
                export_df = export_df.iloc[list(fail_set)]
            else:
                all_idx = range(len(export_df))
                pass_idx = [i for i in all_idx if i not in fail_set]
                export_df = export_df.iloc[pass_idx]

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"

        COLOR_HEADER_BG = "2C3E50"
        header_fill = PatternFill(start_color=COLOR_HEADER_BG, end_color=COLOR_HEADER_BG, fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF", size=11)
        center_align = Alignment(horizontal="center", vertical="center")
        fail_fill = PatternFill(start_color="F5B7B1", end_color="F5B7B1", fill_type="solid")
        fail_font = Font(bold=True, color="FFFFFF")
        thin_border = Border(left=Side(style='thin'), right=Side(style='thin'), top=Side(style='thin'), bottom=Side(style='thin'))

        for col_idx, col_name in enumerate(export_df.columns, 1):
            cell = ws.cell(row=1, column=col_idx, value=col_name)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = center_align
            cell.border = thin_border

        for col_idx, col_name in enumerate(export_df.columns, 1):
            cell = ws.cell(row=2, column=col_idx, value=metadata.get('units', {}).get(col_name, ''))
            cell.alignment = center_align
            cell.border = thin_border

        for col_idx, col_name in enumerate(export_df.columns, 1):
            cell = ws.cell(row=3, column=col_idx, value=metadata.get('mins', {}).get(col_name, ''))
            cell.alignment = center_align
            cell.border = thin_border

        for col_idx, col_name in enumerate(export_df.columns, 1):
            cell = ws.cell(row=4, column=col_idx, value=metadata.get('maxs', {}).get(col_name, ''))
            cell.alignment = center_align
            cell.border = thin_border

        stats_labels = ['Min', 'Avg', 'Max', 'Range', 'STD', 'CPK']
        for r_idx, label in enumerate(stats_labels):
            ws.cell(row=5 + r_idx, column=1, value=label).font = Font(bold=True)

        numeric_cols = [c for c in export_df.columns if export_df[c].dtype in ('int64', 'float64')]
        for col_name in numeric_cols:
            col_idx = list(export_df.columns).index(col_name) + 1
            col_data = pd.to_numeric(export_df[col_name], errors='coerce').dropna()
            if len(col_data) > 0:
                ws.cell(row=5, column=col_idx, value=round(float(col_data.min()), 6))
                ws.cell(row=6, column=col_idx, value=round(float(col_data.mean()), 6))
                ws.cell(row=7, column=col_idx, value=round(float(col_data.max()), 6))
                ws.cell(row=8, column=col_idx, value=round(float(col_data.max() - col_data.min()), 6))
                ws.cell(row=9, column=col_idx, value=round(float(col_data.std()), 6))
                try:
                    min_v = float(metadata['mins'][col_name])
                    max_v = float(metadata['maxs'][col_name])
                    cpk_v = round(min((max_v - col_data.mean()) / (3 * col_data.std()), (col_data.mean() - min_v) / (3 * col_data.std())), 6) if col_data.std() > 0 else 0
                except:
                    cpk_v = 0
                ws.cell(row=10, column=col_idx, value=cpk_v)

        data_start = 12
        fail_indices, _, fail_cells = detect_fail_data(export_df, metadata)
        fail_set = set(fail_indices)
        fail_cells_map = fail_cells
        target_bin = get_bin_column_name(datafile.format_type)

        for r_idx, (_, row) in enumerate(export_df.iterrows()):
            excel_row = data_start + r_idx
            is_fail_row = r_idx in fail_set
            for c_idx, col_name in enumerate(export_df.columns, 1):
                val = row[col_name]
                try:
                    if pd.isna(val):
                        val = ''
                except:
                    pass
                cell = ws.cell(row=excel_row, column=c_idx, value=val)
                cell.border = thin_border
                cell.alignment = center_align
                if is_fail_row and col_name == target_bin:
                    cell.fill = fail_fill
                    cell.font = fail_font
                elif is_fail_row and r_idx in fail_cells_map and col_name in fail_cells_map[r_idx]:
                    cell.fill = fail_fill
                    cell.font = fail_font

        ws.freeze_panes = ws.cell(row=data_start, column=1)

        buffer = io.BytesIO()
        wb.save(buffer)
        buffer.seek(0)

        fname = datafile.filename.rsplit('.', 1)[0]
        response = FileResponse(buffer, as_attachment=True, filename=f'{fname}_analysis.xlsx',
                                content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
        return response

    @action(detail=False, methods=['post'])
    def to_csv(self, request):
        file_id = request.data.get('file_id')
        passfail = request.data.get('passfail', '全部')
        keep_header = request.data.get('keep_header', False)

        datafile = get_object_or_404(DataFile, pk=file_id, owner=request.user)
        parser = get_parser(datafile.format_type)
        df, metadata = parser.parse(datafile.file_path)
        if df is None:
            return Response({'error': 'parse_failed'}, status=400)

        export_df = df.copy()
        if passfail != '全部':
            fail_indices, _, _ = detect_fail_data(export_df, metadata)
            fail_set = set(fail_indices)
            if passfail == 'Pass':
                pass_idx = [i for i in range(len(export_df)) if i not in fail_set]
                export_df = export_df.iloc[pass_idx]
            else:
                export_df = export_df.iloc[list(fail_set)]

        csv_content = export_df.to_csv(index=False)
        fname = datafile.filename.rsplit('.', 1)[0]
        response = FileResponse(io.BytesIO(csv_content.encode('utf-8-sig')), as_attachment=True,
                                filename=f'{fname}_data.csv', content_type='text/csv')
        return response

    @action(detail=False, methods=['post'])
    def sigma_limit(self, request):
        file_id = request.data.get('file_id')
        sigma_level = request.data.get('sigma', 3)
        only_valid = request.data.get('only_valid_limits', False)

        datafile = get_object_or_404(DataFile, pk=file_id, owner=request.user)
        parser = get_parser(datafile.format_type)
        df, metadata = parser.parse(datafile.file_path)
        if df is None:
            return Response({'error': 'parse_failed'}, status=400)

        wb = Workbook()
        ws = wb.active
        ws.title = "TestItem_Limit"
        header_fill = PatternFill(start_color="2C3E50", end_color="2C3E50", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF", size=11)

        for i, h in enumerate(['序号', '测试项', '原LimitL', '原LimitH', f'{sigma_level}σ LimitL', f'{sigma_level}σ LimitH'], 1):
            cell = ws.cell(row=1, column=i, value=h)
            cell.fill = header_fill
            cell.font = header_font

        numeric_cols = [c for c in df.columns if df[c].dtype in ('int64', 'float64')]
        row_idx = 2
        serial = 1
        for param in numeric_cols:
            if param not in metadata.get('mins', {}) or param not in metadata.get('maxs', {}):
                continue
            min_str = str(metadata['mins'][param]).strip()
            max_str = str(metadata['maxs'][param]).strip()
            NON_NUM = ['min', 'max', 'lower limit', 'upper limit', 'n/a', 'na', '-', 'none', '']
            if only_valid and (min_str.lower() in NON_NUM or max_str.lower() in NON_NUM):
                continue

            data_series = get_1d_from(df, param).dropna()
            if len(data_series) == 0:
                continue
            mean_val = float(data_series.mean())
            std_val = float(data_series.std(ddof=0)) if len(data_series) > 1 else 0
            sigma_min = mean_val - sigma_level * std_val
            sigma_max = mean_val + sigma_level * std_val

            ws.cell(row=row_idx, column=1, value=serial)
            ws.cell(row=row_idx, column=2, value=param)
            try:
                ws.cell(row=row_idx, column=3, value=round(float(min_str), 4))
            except:
                ws.cell(row=row_idx, column=3, value='N/A')
            try:
                ws.cell(row=row_idx, column=4, value=round(float(max_str), 4))
            except:
                ws.cell(row=row_idx, column=4, value='N/A')
            ws.cell(row=row_idx, column=5, value=round(sigma_min, 4))
            ws.cell(row=row_idx, column=6, value=round(sigma_max, 4))
            row_idx += 1
            serial += 1

        buffer = io.BytesIO()
        wb.save(buffer)
        buffer.seek(0)
        fname = datafile.filename.rsplit('.', 1)[0]
        return FileResponse(buffer, as_attachment=True, filename=f'{fname}_{sigma_level}sigma_Limit.xlsx',
                            content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

    @action(detail=False, methods=['post'])
    def html_report(self, request):
        file_id = request.data.get('file_id')
        datafile = get_object_or_404(DataFile, pk=file_id, owner=request.user)
        parser = get_parser(datafile.format_type)
        df, metadata = parser.parse(datafile.file_path)
        if df is None:
            return Response({'error': 'parse_failed'}, status=400)

        total_rows = df.shape[0]
        from apps.analysis.services.statistics import calculate_fail_bin_statistics
        bin_stats = calculate_fail_bin_statistics(df, metadata)
        total_pass = 0
        for bv, s in bin_stats.items():
            try:
                if int(float(bv)) == 1:
                    total_pass = s['count']
            except:
                pass
        yield_pct = (total_pass / total_rows * 100) if total_rows > 0 else 0

        html = f"""<!DOCTYPE html><html><head><meta charset="utf-8"><title>ATE Report - {datafile.filename}</title>
<style>body{{font-family:Arial;margin:20px}}h1{{color:#2c3e50}}table{{border-collapse:collapse;width:100%}}th,td{{border:1px solid #ddd;padding:8px;text-align:center}}th{{background:#2c3e50;color:white}}</style></head>
<body><h1>ATE 数据分析报告</h1><p>文件: {datafile.filename} | 格式: {datafile.format_type} | 程序: {datafile.program_name}</p>
<h2>核心指标</h2><table><tr><th>总记录数</th><th>Pass</th><th>Fail</th><th>Yield</th></tr>
<tr><td>{total_rows}</td><td>{total_pass}</td><td>{total_rows - total_pass}</td><td>{yield_pct:.2f}%</td></tr></table></body></html>"""
        return Response(io.BytesIO(html.encode('utf-8')).read(), status=200,
                        content_type='text/html; charset=utf-8',
                        headers={'Content-Disposition': f'attachment; filename="{datafile.filename}_report.html"'})

    @action(detail=False, methods=['post'])
    def batch_charts(self, request):
        file_id = request.data.get('file_id')
        params = request.data.get('params', [])
        fmt = request.data.get('format', 'xlsx')

        if not file_id:
            return Response({'error': 'file_id_required'}, status=400)

        datafile = get_object_or_404(DataFile, pk=file_id, owner=request.user)
        parser = get_parser(datafile.format_type)
        df, metadata = parser.parse(datafile.file_path)
        if df is None:
            return Response({'error': 'parse_failed'}, status=400)

        if not params:
            params = [c for c in df.columns if df[c].dtype in ('int64', 'float64')][:10]

        if fmt == 'pptx':
            return self._batch_charts_pptx(datafile, df, metadata, params)
        else:
            return self._batch_charts_xlsx(datafile, df, metadata, params)

    def _batch_charts_xlsx(self, datafile, df, metadata, params):
        """Export batch stats to Excel."""
        from openpyxl import Workbook
        from openpyxl.styles import PatternFill, Font, Alignment, Border, Side

        wb = Workbook()
        header_fill = PatternFill(start_color="2C3E50", end_color="2C3E50", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF", size=11)
        center_align = Alignment(horizontal="center", vertical="center")
        thin_border = Border(
            left=Side(style='thin'), right=Side(style='thin'),
            top=Side(style='thin'), bottom=Side(style='thin'),
        )

        ws = wb.active
        ws.title = "Summary"

        headers = ['参数', 'Mean', 'STD', 'CPK', 'CPK Level', 'Min', 'Max',
                   'Data Min', 'Data Max', 'Lower Limit', 'Upper Limit',
                   '3σ Min', '3σ Max', '6σ Min', '6σ Max', 'N', 'Unit']
        for i, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=i, value=h)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = center_align
            cell.border = thin_border

        row_idx = 2
        for param in params:
            if param not in df.columns:
                continue
            data_series = get_1d_from(df, param).dropna()
            data_series = data_series[data_series.apply(lambda x: abs(x) < float('inf'))]
            if len(data_series) == 0:
                continue

            stats = compute_range_statistics(data_series, metadata, param)
            cpk_val, cpk_level, cpk_color = compute_cpk(
                stats['mean'], stats['std'], stats['rdl'][0], stats['rdl'][1]
            )

            vals = [
                param,
                round(stats['mean'], 4), round(stats['std'], 4),
                round(cpk_val, 4), cpk_level,
                round(stats['dr'][0], 4), round(stats['dr'][1], 4),
                round(stats['dr'][0], 4), round(stats['dr'][1], 4),
                round(stats['rdl'][0], 4), round(stats['rdl'][1], 4),
                round(stats['s3'][0], 4), round(stats['s3'][1], 4),
                round(stats['s6'][0], 4), round(stats['s6'][1], 4),
                len(data_series), stats['unit'],
            ]
            for i, v in enumerate(vals, 1):
                cell = ws.cell(row=row_idx, column=i, value=v if v is not None else '')
                cell.border = thin_border
                cell.alignment = center_align
            row_idx += 1

        # Auto-fit column widths
        for col_idx in range(1, len(headers) + 1):
            ws.column_dimensions[get_column_letter(col_idx)].width = 16

        buffer = io.BytesIO()
        wb.save(buffer)
        buffer.seek(0)
        fname = datafile.filename.rsplit('.', 1)[0]
        return FileResponse(buffer, as_attachment=True,
                            filename=f'{fname}_batch_stats.xlsx',
                            content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

    def _batch_charts_pptx(self, datafile, df, metadata, params):
        """Export batch charts to PPTX."""
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        # Use blank layout
        blank_layout = prs.slide_layouts[6]  # blank

        # Chinese font support
        plt.rcParams['font.sans-serif'] = ['SimHei', 'Microsoft YaHei', 'DejaVu Sans']
        plt.rcParams['axes.unicode_minus'] = False

        for param in params:
            if param not in df.columns:
                continue
            data_series = get_1d_from(df, param).dropna()
            data_series = data_series[data_series.apply(lambda x: abs(x) < float('inf'))]
            if len(data_series) == 0:
                continue

            stats = compute_range_statistics(data_series, metadata, param)
            cpk_val, cpk_level, _ = compute_cpk(
                stats['mean'], stats['std'], stats['rdl'][0], stats['rdl'][1]
            )

            # Create matplotlib chart
            fig, ax = plt.subplots(figsize=(8, 4.5))
            rdl_min, rdl_max = stats['rdl']
            gap = (rdl_max - rdl_min) / 25 if rdl_max != rdl_min else 0.01
            bin_start = rdl_min - 2.5 * gap
            bins = np.array([bin_start + j * gap for j in range(26)])
            ax.hist(data_series.dropna(), bins=bins, color='#1E88E5', edgecolor='white', alpha=0.85)

            if rdl_min is not None:
                ax.axvline(rdl_min, color='#C62828', linestyle='--', linewidth=2, label='LSL')
            if rdl_max is not None:
                ax.axvline(rdl_max, color='#C62828', linestyle='--', linewidth=2, label='USL')
            if stats.get('s6'):
                ax.axvline(stats['s6'][0], color='#E65100', linestyle=':', linewidth=1.5, label='6σL')
                ax.axvline(stats['s6'][1], color='#E65100', linestyle=':', linewidth=1.5, label='6σU')

            ax.set_title(f'{param}  |  CPK={cpk_val:.4f} ({cpk_level})  |  N={len(data_series)}', fontsize=11)
            ax.set_xlabel(stats.get('unit', ''))
            ax.set_ylabel('Frequency')
            ax.legend(loc='upper right', fontsize=8)
            fig.tight_layout()

            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=120)
            buf.seek(0)
            plt.close(fig)

            slide = prs.slides.add_slide(blank_layout)
            left = Inches(0.5)
            top = Inches(0.5)
            slide.shapes.add_picture(buf, left, top, width=Inches(9), height=Inches(5.5))

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        fname = datafile.filename.rsplit('.', 1)[0]
        return FileResponse(buffer, as_attachment=True,
                            filename=f'{fname}_batch_charts.pptx',
                            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
